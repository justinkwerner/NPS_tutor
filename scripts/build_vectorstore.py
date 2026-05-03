"""Build (or incrementally update) the unified ChromaDB vectorstore.

Usage:
    python scripts/build_vectorstore.py
    python scripts/build_vectorstore.py --data-dir ./data --output ./databases/chroma_unified

Idempotent: chunks already present in the store (matched by ID) are skipped.

Adding a new data source
------------------------
1. Produce a JSON file of chunk dicts (keys: text, source, filename, page, filetype).
   Use scripts/ingest_pdf.py to generate this from PDF/PPTX/DOCX files.
2. Add a DocTypeConfig entry to DEFAULT_DOC_TYPES below.
3. Run this script — new chunks are picked up automatically.

No other code changes are required.
"""

import argparse
import hashlib
import json
import logging
import os
import re
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Callable

from dotenv import load_dotenv
from langchain_chroma import Chroma
from langchain_core.documents import Document
from langchain_experimental.text_splitter import SemanticChunker
from langchain_huggingface import HuggingFaceEmbeddings, HuggingFaceEndpointEmbeddings

try:
    from langchain_nvidia_ai_endpoints import NVIDIAEmbeddings  # type: ignore
except ImportError:
    NVIDIAEmbeddings = None

load_dotenv()

logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")
logger = logging.getLogger(__name__)

os.environ.setdefault("TOKENIZERS_PARALLELISM", "false")

# ── Topic inference helpers ───────────────────────────────────────────────────

MIN_CHUNK_LEN = 20

LAB_TOPIC_MAP = {
    "Lab0": "serial_comm", "ASCII": "serial_comm",
    "Lab1": "serial_comm", "Lab2": "serial_comm", "SerialComm": "serial_comm",
    "Lab3": "comm_protocols", "Polling": "comm_protocols",
    "Lab5": "gpio_switches", "GPIO": "gpio_switches", "Switches": "gpio_switches",
    "Lab6": "design_electronics", "Design_Basic": "design_electronics",
    "Lab7": "spacecraft_ops", "Spacecraft_Ops": "spacecraft_ops", "TIC": "spacecraft_ops",
    "Lab8": "soldering_functional_test", "Soldering": "soldering_functional_test",
    "Functional": "soldering_functional_test",
    "Lab9": "environmental_testing", "Vibe": "environmental_testing",
    "TVAC": "environmental_testing",
    "TCPIP": "networking", "JSON": "networking",
    "Bus-Payload": "bus_payload_comms",
    "Corvus": "spacecraft_bus",
    "EPS": "power_systems",
    "RISP": "risp",
}

LECTURE_TOPIC_MAP = {
    "0": "course_intro", "1": "binary_number_systems", "2": "encoding",
    "3": "cdh_comm_protocols", "4": "ethernet_udp", "5": "basic_circuits",
    "6": "computer_architecture", "7": "power_systems",
    "8": "spacecraft_buses_icds", "9": "testing",
}

MANUAL_TOPIC_KEYWORDS = {
    "trigger": "oscilloscope_trigger", "acquisition": "oscilloscope_acquisition",
    "cursor": "oscilloscope_cursors", "measure": "oscilloscope_measurements",
    "math": "oscilloscope_math", "channel": "oscilloscope_channels",
    "probe": "oscilloscope_probes", "display": "oscilloscope_display",
    "ethernet": "oscilloscope_ethernet", "print": "oscilloscope_printing",
    "save": "oscilloscope_save_recall", "recall": "oscilloscope_save_recall",
    "setup": "oscilloscope_setup", "fft": "oscilloscope_fft",
    "decode": "oscilloscope_bus_decode", "bus": "oscilloscope_bus_decode",
    "waveform": "oscilloscope_waveform", "zoom": "oscilloscope_zoom",
    "video": "oscilloscope_video", "safety": "oscilloscope_safety",
    "specification": "oscilloscope_specs", "bandwidth": "oscilloscope_specs",
}

# TODO (Phase 8): Add SATOPS_TOPIC_KEYWORDS here once satellite ops manual is ingested.
# Pattern: {"keyword": "satops_<subtopic>", ...}
# SATOPS_TOPIC_KEYWORDS = {}


def _extract_lab_number(filename: str) -> int | None:
    m = re.search(r"[Ll]ab\s*(\d+)", filename)
    return int(m.group(1)) if m else None


def _extract_lecture_number(filename: str) -> int | None:
    m = re.search(r"SS3861\s*-\s*(\d+)\s*-", filename)
    return int(m.group(1)) if m else None


def _infer_lab_topic(text: str, filename: str) -> str:
    for pattern, topic in LAB_TOPIC_MAP.items():
        if pattern.lower() in filename.lower():
            return topic
    return "lab_general"


def _infer_lecture_topic(text: str, filename: str) -> str:
    num = _extract_lecture_number(filename)
    if num is not None:
        return LECTURE_TOPIC_MAP.get(str(num), "lecture_general")
    if "polling" in filename.lower() or "blocking" in filename.lower():
        return "comm_protocols"
    if "Lectures.pdf" in filename:
        return "lecture_notes_compiled"
    return "lecture_general"


def _infer_manual_topic(text: str, filename: str) -> str:
    text_lower = text.lower()
    best_topic, best_count = "oscilloscope_general", 0
    for keyword, topic in MANUAL_TOPIC_KEYWORDS.items():
        count = text_lower.count(keyword)
        if count > best_count:
            best_count, best_topic = count, topic
    return best_topic


# TODO (Phase 8): Implement _infer_satops_topic once SATOPS_TOPIC_KEYWORDS is populated.
def _infer_satops_topic(text: str, filename: str) -> str:
    return "satops_general"


# ── DocTypeConfig ─────────────────────────────────────────────────────────────

@dataclass
class DocTypeConfig:
    """Describes one category of source documents and how to enrich its chunks.

    Args:
        doc_type:          Value stored in metadata["doc_type"] (e.g. "lab").
        json_filename:     JSON file name relative to data_dir (e.g. "labs.json").
        topic_fn:          Callable(text, filename) -> topic string.
        lab_number_fn:     Callable(filename) -> int | None. Defaults to always None.
        lecture_number_fn: Callable(filename) -> int | None. Defaults to always None.
        required:          If False, a missing JSON file is a warning, not an error.
                           Use False for data sources that are planned but not yet ingested.
    """
    doc_type: str
    json_filename: str
    topic_fn: Callable[[str, str], str]
    lab_number_fn: Callable[[str], int | None] = field(default=lambda fn: None)
    lecture_number_fn: Callable[[str], int | None] = field(default=lambda fn: None)
    required: bool = True


# ── Default doc type registry ─────────────────────────────────────────────────
# To add a new data source: append a DocTypeConfig here. That's it.

DEFAULT_DOC_TYPES: list[DocTypeConfig] = [
    DocTypeConfig(
        doc_type="lab",
        json_filename="labs.json",
        topic_fn=_infer_lab_topic,
        lab_number_fn=_extract_lab_number,
    ),
    DocTypeConfig(
        doc_type="lecture",
        json_filename="lectures.json",
        topic_fn=_infer_lecture_topic,
        lecture_number_fn=_extract_lecture_number,
    ),
    DocTypeConfig(
        doc_type="manual",
        json_filename="manual.json",
        topic_fn=_infer_manual_topic,
    ),
    # ── Phase 8: Satellite Operations Manual ──────────────────────────────────
    # TODO: Run scripts/ingest_pdf.py on the satellite ops PDF to produce satops.json,
    #       then set required=True here.
    DocTypeConfig(
        doc_type="satops",
        json_filename="satops.json",
        topic_fn=_infer_satops_topic,
        required=False,
    ),
]


# ── Generic chunk enrichment ──────────────────────────────────────────────────

def _chunk_id(chunk: dict, doc_type: str) -> str:
    """Stable ID derived from source path + content hash."""
    source = chunk.get("source", chunk.get("filename", "unknown"))
    content_hash = hashlib.md5(chunk.get("text", "").encode()).hexdigest()[:12]
    return f"{doc_type}_{source}_{content_hash}"


def _enrich(chunks: list, cfg: DocTypeConfig) -> list[Document]:
    """Convert raw chunk dicts into enriched LangChain Documents using cfg."""
    docs = []
    for chunk in chunks:
        text = chunk.get("text", "").strip()
        if len(text) < MIN_CHUNK_LEN:
            continue
        filename = chunk.get("filename", "unknown")
        docs.append(Document(
            page_content=text,
            metadata={
                "source": chunk.get("source", "unknown"),
                "filename": filename,
                "page": chunk.get("page", 0),
                "filetype": chunk.get("filetype", "unknown"),
                "doc_type": cfg.doc_type,
                "lab_number": cfg.lab_number_fn(filename),
                "lecture_number": cfg.lecture_number_fn(filename),
                "topic": cfg.topic_fn(text, filename),
                "chunk_id": _chunk_id(chunk, cfg.doc_type),
            }
        ))
    return docs


# ── Load & enrich ─────────────────────────────────────────────────────────────

def load_and_enrich(
    data_dir: str = ".",
    configs: list[DocTypeConfig] | None = None,
) -> list[Document]:
    """Load JSON chunk files and return enriched Documents.

    Args:
        data_dir: Directory containing the JSON files listed in configs.
        configs:  List of DocTypeConfig entries. Defaults to DEFAULT_DOC_TYPES.

    Returns:
        Flat list of enriched Documents across all configured doc types.
    """
    if configs is None:
        configs = DEFAULT_DOC_TYPES

    data = Path(data_dir)
    all_docs: list[Document] = []
    counts: dict[str, int] = {}

    for cfg in configs:
        json_path = data / cfg.json_filename
        if not json_path.exists():
            if cfg.required:
                raise FileNotFoundError(
                    f"Required data file not found: {json_path}\n"
                    f"Run scripts/ingest_pdf.py to generate it."
                )
            logger.warning(
                "Skipping optional doc type '%s' — file not found: %s",
                cfg.doc_type, json_path,
            )
            continue

        with open(json_path, encoding="utf-8") as f:
            chunks = json.load(f)

        docs = _enrich(chunks, cfg)
        counts[cfg.doc_type] = len(docs)
        all_docs.extend(docs)
        logger.info("Loaded '%s': %d chunks from %s", cfg.doc_type, len(docs), json_path.name)

    logger.info("Total enriched chunks: %d across %d doc types", len(all_docs), len(counts))

    topic_counts = Counter(d.metadata["topic"] for d in all_docs)
    logger.info("Topic distribution (%d topics):", len(topic_counts))
    for topic, count in topic_counts.most_common():
        logger.info("  %-40s %4d", topic, count)

    return all_docs


# ── Embedding setup ───────────────────────────────────────────────────────────

def _setup_embeddings() -> object:
    embed_model = os.getenv("NVIDIA_EMBED_MODEL", "nvidia/nv-embed-v1")
    embed_server_url = os.getenv(
        "EMBEDDING_SERVER_URL",
        "http://trac-malenia.ern.nps.edu:8080/gpu5/embed",
    )
    fallback_model = "all-MiniLM-L6-v2"
    nvidia_key = os.getenv("ALTERNATE_EMBEDDING_MODEL_KEY")
    nvidia_url = os.getenv("NIM_BASE_URL", "https://integrate.api.nvidia.com/v1")

    # 1) NVIDIA hosted embeddings
    if NVIDIAEmbeddings is not None and nvidia_key:
        try:
            hfe = NVIDIAEmbeddings(model=embed_model, api_key=nvidia_key, base_url=nvidia_url)
            hfe.embed_query("test")
            logger.info("Embedding function set from NVIDIA API: %s", embed_model)
            return hfe
        except Exception as e:
            logger.warning("NVIDIA embeddings failed (%s) — trying HF endpoint", e)

    # 2) NPS HuggingFace endpoint server
    try:
        hfe = HuggingFaceEndpointEmbeddings(model=embed_server_url)
        hfe.embed_query("test")
        logger.info("Embedding function set from HF endpoint server: %s", embed_server_url)
        return hfe
    except Exception as e:
        logger.info("HF endpoint unreachable (%s) — using local embeddings", e)

    # 3) Local fallback
    hfe = HuggingFaceEmbeddings(
        model_name=fallback_model,
        model_kwargs={"device": "cpu"},
        encode_kwargs={"normalize_embeddings": False},
    )
    logger.info("Embedding function set locally with %s", fallback_model)
    return hfe


# ── Vectorstore build ─────────────────────────────────────────────────────────

def _apply_semantic_chunking(docs: list[Document], embeddings) -> list[Document]:
    """Apply SemanticChunker as a second pass, preserving metadata from the source doc."""
    chunker = SemanticChunker(embeddings=embeddings, breakpoint_threshold_type="percentile")
    split_docs = chunker.transform_documents(docs)
    for doc in split_docs:
        source = doc.metadata.get("source", doc.metadata.get("filename", "unknown"))
        doc_type = doc.metadata.get("doc_type", "unknown")
        content_hash = hashlib.md5(doc.page_content.encode()).hexdigest()[:12]
        doc.metadata["chunk_id"] = f"{doc_type}_{source}_{content_hash}"
    return split_docs


def build_vectorstore(
    data_dir: str = ".",
    output_dir: str = "./databases/chroma_unified",
    collection_name: str = "unified_course",
    configs: list[DocTypeConfig] | None = None,
) -> Chroma:
    """Build or incrementally update the unified ChromaDB vectorstore.

    Args:
        data_dir:        Directory containing the JSON files.
        output_dir:      Chroma persist directory.
        collection_name: Chroma collection name.
        configs:         Doc type configs. Defaults to DEFAULT_DOC_TYPES.

    Returns:
        The loaded/updated Chroma vectorstore instance.
    """
    all_docs = load_and_enrich(data_dir=data_dir, configs=configs)
    hfe = _setup_embeddings()

    # Apply semantic chunking per doc type.
    by_type: dict[str, list[Document]] = {}
    for doc in all_docs:
        dt = doc.metadata["doc_type"]
        by_type.setdefault(dt, []).append(doc)

    chunked: list[Document] = []
    for dt, docs in by_type.items():
        logger.info("Applying semantic chunking to '%s' (%d docs)...", dt, len(docs))
        result = _apply_semantic_chunking(docs, hfe)
        logger.info("  → %d chunks after semantic chunking", len(result))
        chunked.extend(result)

    logger.info("Total after semantic chunking: %d", len(chunked))

    vectorstore = Chroma(
        collection_name=collection_name,
        embedding_function=hfe,
        persist_directory=output_dir,
    )

    # Idempotency: skip chunks whose ID is already in the store.
    existing_ids = set(vectorstore.get()["ids"])
    logger.info("Existing chunks in store: %d", len(existing_ids))

    seen_ids: set[str] = set(existing_ids)
    new_docs = []
    for d in chunked:
        cid = d.metadata["chunk_id"]
        if cid not in seen_ids:
            seen_ids.add(cid)
            new_docs.append(d)
    logger.info("New chunks to add: %d", len(new_docs))

    if not new_docs:
        logger.info("Vectorstore is already up to date — nothing to add.")
        return vectorstore

    ids = [d.metadata["chunk_id"] for d in new_docs]
    batch_size = 64
    for i in range(0, len(new_docs), batch_size):
        batch = new_docs[i:i + batch_size]
        vectorstore.add_documents(documents=batch, ids=ids[i:i + batch_size])
        logger.info("  Added batch %d–%d", i + 1, i + len(batch))

    total = len(vectorstore.get()["ids"])
    logger.info("Done. Total chunks in store: %d", total)
    return vectorstore


# ── Single-file ingestion (for instructor uploads) ────────────────────────────

def ingest_file_to_vectorstore(
    file_path: str,
    doc_type: str = "upload",
    output_dir: str = "./databases/chroma_unified",
    collection_name: str = "unified_course",
    original_filename: str | None = None,
) -> int:
    """Ingest a single file into the existing vectorstore.

    Uses Docling for conversion, then the same enrich → semantic-chunk →
    idempotent-add pipeline as build_vectorstore().

    Args:
        file_path:       Absolute path to the file (PDF, PPTX, DOCX, …).
        doc_type:        Metadata label stored in Chroma (default "upload").
        output_dir:      Chroma persist directory.
        collection_name: Chroma collection name.

    Returns:
        Number of new chunks added to the store.
    """
    path = Path(file_path)
    display_name = original_filename or path.name
    logger.info("Ingesting uploaded file: %s (doc_type=%s)", display_name, doc_type)

    raw_chunks: list[dict] = []

    if path.suffix.lower() == ".md":
        # Markdown: read directly and split — skip Docling entirely.
        import re
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        text = path.read_text(encoding="utf-8")
        # Strip YAML frontmatter (--- ... ---) if present.
        text = re.sub(r"^---\s*\n.*?\n---\s*\n", "", text, flags=re.DOTALL)

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=512,
            chunk_overlap=64,
            separators=["\n## ", "\n### ", "\n#### ", "\n\n", "\n", " "],
        )
        for chunk_text in splitter.split_text(text):
            if chunk_text.strip():
                raw_chunks.append({
                    "text": chunk_text.strip(),
                    "source": display_name,
                    "filename": display_name,
                    "page": 0,
                    "filetype": ".md",
                })
        logger.info("  Markdown split produced %d raw chunks", len(raw_chunks))

    elif path.suffix.lower() == ".py":
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        text = path.read_text(encoding="utf-8")
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=512,
            chunk_overlap=64,
            separators=["\nclass ", "\ndef ", "\n\n", "\n", " "],
        )
        for chunk_text in splitter.split_text(text):
            if chunk_text.strip():
                raw_chunks.append({
                    "text": chunk_text.strip(),
                    "source": display_name,
                    "filename": display_name,
                    "page": 0,
                    "filetype": ".py",
                })
        logger.info("  Python split produced %d raw chunks", len(raw_chunks))

    else:
        # Lightweight extraction — no ML models required.
        # Uses pypdf (PDF), python-pptx (PPTX), python-docx (DOCX).
        from langchain_text_splitters import RecursiveCharacterTextSplitter

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=512,
            chunk_overlap=64,
            separators=["\n\n", "\n", " "],
        )

        suffix = path.suffix.lower()

        if suffix == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            pages = [(page.extract_text() or "", i + 1) for i, page in enumerate(reader.pages)]
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            pages = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts = [
                    shape.text_frame.text
                    for shape in slide.shapes
                    if shape.has_text_frame
                ]
                pages.append(("\n".join(texts), slide_num))
        elif suffix == ".docx":
            from docx import Document as DocxDocument
            doc_obj = DocxDocument(str(path))
            full_text = "\n".join(p.text for p in doc_obj.paragraphs)
            pages = [(full_text, 0)]
        else:
            raise ValueError(f"Unsupported file type for lightweight ingestion: {suffix}")

        for page_text, page_num in pages:
            if not page_text.strip():
                continue
            for chunk_text in splitter.split_text(page_text):
                if chunk_text.strip():
                    raw_chunks.append({
                        "text": chunk_text.strip(),
                        "source": display_name,
                        "filename": display_name,
                        "page": page_num,
                        "filetype": suffix,
                    })

        logger.info("  Lightweight extraction produced %d raw chunks", len(raw_chunks))

    cfg = DocTypeConfig(
        doc_type=doc_type,
        json_filename="",
        topic_fn=lambda text, filename: "uploaded_content",
    )
    docs = _enrich(raw_chunks, cfg)

    hfe = _setup_embeddings()
    chunked = _apply_semantic_chunking(docs, hfe)
    logger.info("  %d chunks after semantic chunking", len(chunked))

    vectorstore = Chroma(
        collection_name=collection_name,
        embedding_function=hfe,
        persist_directory=output_dir,
    )

    existing_ids = set(vectorstore.get()["ids"])
    new_docs = [d for d in chunked if d.metadata["chunk_id"] not in existing_ids]
    logger.info("  New chunks to add: %d (skipping %d duplicates)", len(new_docs), len(chunked) - len(new_docs))

    if not new_docs:
        return 0

    ids = [d.metadata["chunk_id"] for d in new_docs]
    batch_size = 64
    for i in range(0, len(new_docs), batch_size):
        batch = new_docs[i:i + batch_size]
        vectorstore.add_documents(documents=batch, ids=ids[i:i + batch_size])

    logger.info("  Done. Added %d chunks for '%s'", len(new_docs), display_name)
    return len(new_docs)


# ── CLI ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Build the unified ChromaDB vectorstore.")
    parser.add_argument(
        "--data-dir", default=str(Path(__file__).parent.parent),
        help="Directory containing JSON data files (default: project root)"
    )
    parser.add_argument(
        "--output", default=os.getenv("CHROMA_DB_PATH", "./databases/chroma_unified"),
        help="Chroma persist directory (default: ./databases/chroma_unified)"
    )
    args = parser.parse_args()

    build_vectorstore(data_dir=args.data_dir, output_dir=args.output)
